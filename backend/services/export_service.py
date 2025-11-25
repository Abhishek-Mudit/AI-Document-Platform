"""
Document export service for generating .docx and .pptx files
"""
from docx import Document
from docx.shared import Pt, Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from typing import List, Dict
import json
from io import BytesIO


class ExportService:
    """Service for exporting documents to .docx and .pptx formats"""
    
    @staticmethod
    def export_docx(topic: str, sections: List[Dict[str, str]]) -> BytesIO:
        """
        Export content to a Word document
        
        Args:
            topic: Document topic/title
            sections: List of dicts with 'title' and 'content' keys
        
        Returns:
            BytesIO object containing the .docx file
        """
        doc = Document()
        
        # Add title
        title = doc.add_heading(topic, 0)
        title.alignment = 1  # Center alignment
        
        # Add sections
        for section in sections:
            # Section heading
            doc.add_heading(section['title'], level=1)
            
            # Section content
            content = section.get('content', '')
            if content:
                # Split by newlines and add paragraphs
                paragraphs = content.split('\n')
                for para_text in paragraphs:
                    if para_text.strip():
                        para = doc.add_paragraph(para_text.strip())
                        para.style = 'Normal'
        
        # Save to BytesIO
        file_stream = BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        return file_stream
    
    @staticmethod
    def export_pptx(topic: str, slides: List[Dict[str, str]]) -> BytesIO:
        """
        Export content to a PowerPoint presentation
        
        Args:
            topic: Presentation topic/title
            slides: List of dicts with 'title' and 'content' keys
        
        Returns:
            BytesIO object containing the .pptx file
        """
        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = topic
        
        # Content slides
        for slide_data in slides:
            # Use title and content layout
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data['title']
            
            # Set content
            content = slide_data.get('content', '')
            if content and len(slide.shapes) > 1:
                text_frame = slide.shapes[1].text_frame
                text_frame.clear()
                
                # Parse content into bullet points
                lines = content.split('\n')
                for i, line in enumerate(lines):
                    if line.strip():
                        if i == 0:
                            text_frame.text = line.strip()
                        else:
                            p = text_frame.add_paragraph()
                            p.text = line.strip()
                            p.level = 0
        
        # Save to BytesIO
        file_stream = BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        return file_stream


# Global export service instance
export_service = ExportService()
